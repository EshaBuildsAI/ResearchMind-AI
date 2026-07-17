"""
generators.py
Turns AI-generated content into downloadable files
(DOCX reports, PPTX presentations, TXT flashcard decks).
"""

import io
import os
import re
from datetime import datetime

import docx
from pptx import Presentation
from pptx.util import Inches, Pt

from constants import EXPORT_DIR
from utils import ensure_dir, safe_filename


def _export_path(base_name: str, extension: str) -> str:
    ensure_dir(EXPORT_DIR)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = safe_filename(f"{base_name}_{timestamp}.{extension}")
    return os.path.join(EXPORT_DIR, filename)


def export_summary_to_docx(summary_text: str, source_filename: str) -> str:
    """Save a summary as a Word document. Returns the saved file path."""
    document = docx.Document()
    document.add_heading("Document Summary", level=1)
    document.add_paragraph(f"Source: {source_filename}")
    document.add_paragraph("")

    for line in summary_text.split("\n"):
        if line.strip():
            document.add_paragraph(line.strip(), style="List Bullet")

    path = _export_path("summary", "docx")
    document.save(path)
    return path


def export_research_gap_to_docx(analysis_text: str, source_filename: str) -> str:
    """Save a research gap analysis as a Word document."""
    document = docx.Document()
    document.add_heading("Research Gap Analysis", level=1)
    document.add_paragraph(f"Source: {source_filename}")
    document.add_paragraph("")

    for line in analysis_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.endswith(":") and len(line.split()) <= 6:
            document.add_heading(line, level=2)
        else:
            document.add_paragraph(line)

    path = _export_path("research_gap", "docx")
    document.save(path)
    return path


def export_literature_review_to_docx(review_text: str, source_filename: str) -> str:
    """Save a literature review as a Word document."""
    document = docx.Document()
    document.add_heading("Literature Review", level=1)
    document.add_paragraph(f"Source: {source_filename}")
    document.add_paragraph("")

    for line in review_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if re.match(r"^\d+\.\s", line) or line.endswith(":"):
            document.add_heading(line, level=2)
        else:
            document.add_paragraph(line)

    path = _export_path("literature_review", "docx")
    document.save(path)
    return path


def export_workspace_bundle(source_filename: str, sections: dict) -> str:
    """
    Combine every generated section (summary, literature review, research gap, etc.)
    for the active document into ONE downloadable report.
    `sections` example: {"Summary": "...", "Research Gap Analysis": "...", ...}
    Only non-empty sections are included.
    """
    document = docx.Document()
    document.add_heading(f"Research Workspace Report", level=1)
    document.add_paragraph(f"Source document: {source_filename}")
    document.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    for title, content in sections.items():
        if not content:
            continue
        document.add_page_break()
        document.add_heading(title, level=2)
        for line in str(content).split("\n"):
            if line.strip():
                document.add_paragraph(line.strip())

    path = _export_path("workspace_report", "docx")
    document.save(path)
    return path
    """Save a quiz (with answer key) as a Word document."""
    document = docx.Document()
    document.add_heading("Quiz", level=1)
    document.add_paragraph(f"Source: {source_filename}")
    document.add_paragraph("")

    for i, q in enumerate(quiz, start=1):
        document.add_paragraph(f"Q{i}: {q['question']}", style="List Number")
        for letter, option in q["options"].items():
            document.add_paragraph(f"   {letter}) {option}")

    document.add_page_break()
    document.add_heading("Answer Key", level=2)
    for i, q in enumerate(quiz, start=1):
        document.add_paragraph(f"Q{i}: {q['correct']}")

    path = _export_path("quiz", "docx")
    document.save(path)
    return path


def export_flashcards_to_docx(cards: list, source_filename: str) -> str:
    """Save flashcards as a simple two-column Word table."""
    document = docx.Document()
    document.add_heading("Flashcards", level=1)
    document.add_paragraph(f"Source: {source_filename}")

    table = document.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    header = table.rows[0].cells
    header[0].text = "Front"
    header[1].text = "Back"

    for card in cards:
        row = table.add_row().cells
        row[0].text = card["front"]
        row[1].text = card["back"]

    path = _export_path("flashcards", "docx")
    document.save(path)
    return path


def export_presentation_to_pptx(slides: list, title: str = "Research Presentation") -> str:
    """Save a slide outline as an actual .pptx presentation."""
    presentation = Presentation()

    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by ResearchMind AI"

    bullet_layout = presentation.slide_layouts[1]
    for slide_data in slides:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data["title"]
        body = slide.placeholders[1].text_frame
        body.clear()

        for i, bullet in enumerate(slide_data["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)

    path = _export_path("presentation", "pptx")
    presentation.save(path)
    return path
